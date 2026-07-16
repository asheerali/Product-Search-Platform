"""
PPTX parser — extracts:
  • slide text (title + body)
  • table content from slides
  • embedded images saved to pics/
"""
import logging
from dataclasses import dataclass, field
from pathlib import Path

logger = logging.getLogger(__name__)


@dataclass
class ParsedSlide:
    slide_number: int
    title: str
    text: str
    tables: list[list[dict]] = field(default_factory=list)
    image_paths: list[str] = field(default_factory=list)


@dataclass
class PPTXParseResult:
    document_id: str
    filename: str
    slides: list[ParsedSlide] = field(default_factory=list)

    @property
    def full_text(self) -> str:
        return "\n\n".join(
            f"Slide {s.slide_number}: {s.title}\n{s.text}" for s in self.slides
        )

    @property
    def all_tables(self) -> list[list[dict]]:
        tables = []
        for s in self.slides:
            tables.extend(s.tables)
        return tables

    @property
    def all_image_paths(self) -> list[str]:
        paths = []
        for s in self.slides:
            paths.extend(s.image_paths)
        return paths


def parse_pptx(file_path: str, document_id: str, pics_base_dir: str) -> PPTXParseResult:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Pt

    from app.services.ingestion.image_extractor import save_extracted_image

    path = Path(file_path)
    result = PPTXParseResult(document_id=document_id, filename=path.name)

    try:
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, start=1):
            title_text = ""
            body_parts = []
            tables = []
            image_paths = []

            for shape in slide.shapes:
                # --- Title ---
                if shape.has_text_frame and shape.shape_type != MSO_SHAPE_TYPE.TABLE:
                    text = "\n".join(
                        para.text for para in shape.text_frame.paragraphs if para.text.strip()
                    )
                    if shape.name and "title" in shape.name.lower():
                        title_text = text
                    else:
                        body_parts.append(text)

                # --- Tables ---
                if shape.has_table:
                    table = shape.table
                    if table.rows:
                        headers = [
                            cell.text.strip() or f"col_{ci}"
                            for ci, cell in enumerate(table.rows[0].cells)
                        ]
                        rows = []
                        for row in table.rows[1:]:
                            row_dict = {headers[ci]: cell.text.strip() for ci, cell in enumerate(row.cells)}
                            rows.append(row_dict)
                        if rows:
                            tables.append(rows)

                # --- Images ---
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        source_ref = f"slide_{slide_num}_shape_{shape.shape_id}"
                        saved_path = save_extracted_image(
                            image_bytes=image.blob,
                            ext=image.ext,
                            document_id=document_id,
                            source_ref=source_ref,
                            pics_base_dir=pics_base_dir,
                        )
                        if saved_path:
                            image_paths.append(saved_path)
                    except Exception as e:
                        logger.warning("Could not extract image from slide %d: %s", slide_num, e)

            result.slides.append(
                ParsedSlide(
                    slide_number=slide_num,
                    title=title_text,
                    text="\n".join(body_parts),
                    tables=tables,
                    image_paths=image_paths,
                )
            )
    except Exception as e:
        logger.error("PPTX parsing failed for %s: %s", file_path, e)

    return result
